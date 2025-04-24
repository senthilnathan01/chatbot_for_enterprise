# --- START OF FILE file_parsers.py ---
"""Functions to parse different file types.
Each parser should return a list of tuples:
[(text_segment_1, metadata_1), (text_segment_2, metadata_2), ...]
where metadata includes at least 'source' and potentially 'page_number' etc.
Also returns the list of found URLs.

Data files (CSV/XLSX/JSON) are handled differently: they are imported into
a chat-specific database via DataImporter_Gemini, and a marker segment
is returned for ChromaDB indexing.
"""

import streamlit as st
import io
import fitz # PyMuPDF
import docx
from pptx import Presentation
import pandas as pd
import json
import os
from pathlib import Path

# Import specific analysis tool classes
from DataImporter_Gemini import DataImporter_Gemini

from utils import log_message, find_urls
from image_processor import get_gemini_vision_description_ocr
from config import CHAT_DB_DIRECTORY, SUPPORTED_DATA_TYPES, SUPPORTED_TEXT_TYPES

# --- PDF Parser ---
# ... (keep as is) ...
def parse_pdf(file_content, filename):
    parsed_segments = [] # List of (text_segment, metadata)
    urls = set()
    try:
        doc = fitz.open(stream=file_content, filetype="pdf")
        log_message(f"Processing PDF '{filename}': {len(doc)} pages.", "info")
        for page_num, page in enumerate(doc):
            page_num_actual = page_num + 1 # 1-based page number
            page_text = ""
            page_meta = {"source": filename, "page_number": page_num_actual}
            ocr_text_content = "" # Store OCR result separately

            try:
                page_text = page.get_text("text", sort=True) # Get text sorted by position
                if page_text: # Don't add empty pages unless OCR finds something
                    urls.update(find_urls(page_text))
                    parsed_segments.append((page_text, page_meta.copy()))
                else:
                    log_message(f"No extractable text found on page {page_num_actual} of {filename}. Checking for images/OCR.", "debug")

                # OCR Fallback check (can run even if some text was found, for mixed pages)
                ocr_needed = len(page_text.strip()) < 100 # Heuristic: try OCR if very little text
                if not page_text or ocr_needed:
                    log_message(f"Attempting OCR for page {page_num_actual} in '{filename}'.", "info")
                    try:
                        pix = page.get_pixmap(dpi=150)
                        img_bytes = pix.tobytes("png")
                        ocr_text_content = get_gemini_vision_description_ocr(img_bytes, f"{filename}_page{page_num_actual}_ocr")
                        if ocr_text_content and "No description or text extracted" not in ocr_text_content:
                             ocr_meta = page_meta.copy()
                             ocr_meta["content_type"] = "ocr_page"
                             parsed_segments.append((f"[OCR Result for Page {page_num_actual}]\n{ocr_text_content}", ocr_meta))
                        else:
                             log_message(f"OCR found no significant text on page {page_num_actual}.", "debug")
                    except Exception as ocr_err:
                        log_message(f"Error during OCR fallback for {filename} page {page_num_actual}: {ocr_err}", "warning")

                # Extract Embedded Images (if any)
                try:
                    image_list = page.get_images(full=True)
                    for img_index, img_info in enumerate(image_list):
                        xref = img_info[0]
                        base_image = doc.extract_image(xref)
                        if base_image:
                            image_bytes = base_image["image"]
                            img_meta_filename = f"{filename}_page{page_num_actual}_img{img_index}"
                            img_desc_ocr = get_gemini_vision_description_ocr(image_bytes, img_meta_filename)
                            if img_desc_ocr and "No description or text extracted" not in img_desc_ocr:
                                 img_meta = page_meta.copy()
                                 img_meta["content_type"] = "embedded_image"
                                 img_meta["image_ref"] = img_index
                                 parsed_segments.append((f"[Analysis of Embedded Image {img_index} on Page {page_num_actual}]\n{img_desc_ocr}", img_meta))
                        else:
                            log_message(f"Could not extract image xref {xref} on page {page_num_actual}.", "warning")
                except Exception as img_err:
                    log_message(f"Error processing embedded images on page {page_num_actual}: {img_err}", "warning")

            except Exception as page_err:
                 log_message(f"Error processing page {page_num_actual} of {filename}: {page_err}", "warning")

        log_message(f"Finished processing PDF '{filename}'. Found URLs: {len(urls)}")
    except Exception as e:
        log_message(f"Critical error processing PDF '{filename}': {e}", "error")
        # Return empty lists on critical failure
        return [], []
    return parsed_segments, list(urls)


# --- DOCX Parser ---
# ... (keep as is) ...
def parse_docx(file_content, filename):
    parsed_segments = []
    urls = set()
    base_meta = {"source": filename, "content_type": "text"}
    img_meta = {"source": filename, "content_type": "embedded_image"}
    text_content = ""
    try:
        doc_stream = io.BytesIO(file_content)
        doc = docx.Document(doc_stream)
        log_message(f"Processing DOCX '{filename}'.", "info")

        # Extract text from paragraphs
        para_texts = []
        for para in doc.paragraphs:
            para_text = para.text
            para_texts.append(para_text)
            urls.update(find_urls(para_text))
        text_content = '\n'.join(para_texts) # Join with single newline might be better for context
        if text_content:
             parsed_segments.append((text_content, base_meta.copy()))


        # Extract text from tables (basic)
        table_texts = []
        for i, table in enumerate(doc.tables):
             table_str = f"\n[Table {i+1} Content:]\n"
             try:
                  for row in table.rows:
                       row_text = [cell.text for cell in row.cells]
                       table_str += "| " + " | ".join(row_text) + " |\n"
                  table_texts.append(table_str)
             except Exception as table_err:
                  log_message(f"Error reading table {i+1} in {filename}: {table_err}", "warning")
        if table_texts:
            table_meta = base_meta.copy()
            table_meta["content_type"] = "table_text"
            parsed_segments.append(("\n".join(table_texts), table_meta))


        # Handle hyperlinks explicitly
        for rel in doc.part.rels.values():
            if rel.reltype == docx.opc.constants.RELATIONSHIP_TYPE.HYPERLINK and rel.is_external:
                urls.add(rel.target_ref)

        # Image extraction (simplified)
        image_parts = {}
        for rel_id, rel in doc.part.rels.items():
             # Use endswith for broader image format compatibility if needed
             if rel.target_ref.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff')):
                 try:
                      image_parts[rel_id] = rel.target_part.blob
                 except Exception as blob_err:
                      log_message(f"Could not access blob for image rel {rel_id} in {filename}: {blob_err}", "debug")

        img_count = 0
        processed_rel_ids = set()
        # Iterate through shapes that could contain images
        # This includes inline shapes and potentially shapes within text boxes or groups if needed (more complex)
        all_shapes = list(doc.inline_shapes)
        # Add other shape containers if necessary:
        # for section in doc.sections:
        #     for shape in section.header.shapes: all_shapes.append(shape)
        #     for shape in section.footer.shapes: all_shapes.append(shape)
        # for para in doc.paragraphs:
        #      for run in para.runs:
        #          # Check runs for shapes if they can contain them in python-docx
        #          pass

        for shape in all_shapes: # Primarily doc.inline_shapes
             try:
                 # Simplified check based on common structure, might need refinement
                 if hasattr(shape, '_inline') and hasattr(shape._inline, 'graphic') and hasattr(shape._inline.graphic, 'graphicData'):
                     # Use namespaces correctly
                     pic_ns = "{http://schemas.openxmlformats.org/drawingml/2006/picture}"
                     main_ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
                     rel_ns = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"

                     pic = shape._inline.graphic.graphicData.find(f".//{pic_ns}pic")
                     if pic is not None:
                         blip = pic.find(f".//{main_ns}blip")
                         if blip is not None:
                             embed_id = blip.get(f"{rel_ns}embed")
                             if embed_id and embed_id in image_parts and embed_id not in processed_rel_ids:
                                 image_bytes = image_parts[embed_id]
                                 img_meta_filename = f"{filename}_img{img_count}"
                                 img_desc_ocr = get_gemini_vision_description_ocr(image_bytes, img_meta_filename)
                                 if img_desc_ocr and "No description or text extracted" not in img_desc_ocr:
                                      current_img_meta = img_meta.copy()
                                      current_img_meta["image_ref"] = f"inline_{img_count}_{embed_id}"
                                      parsed_segments.append((f"[Analysis of Embedded Image {img_count}]\n{img_desc_ocr}", current_img_meta))
                                      img_count += 1
                                      processed_rel_ids.add(embed_id) # Mark as processed
             except Exception as img_ex:
                 log_message(f"Minor error processing inline shape in {filename}: {img_ex}", "debug")


        log_message(f"Finished processing DOCX '{filename}'. Found URLs: {len(urls)}, Images processed: {img_count}")
    except Exception as e:
        log_message(f"Critical error processing DOCX '{filename}': {e}", "error")
        return [], []
    return parsed_segments, list(urls)

# --- PPTX Parser ---
# ... (keep as is) ...
def parse_pptx(file_content, filename):
    parsed_segments = []
    urls = set()
    try:
        ppt_stream = io.BytesIO(file_content)
        prs = Presentation(ppt_stream)
        log_message(f"Processing PPTX '{filename}': {len(prs.slides)} slides.", "info")

        img_count = 0 # Keep track of images across slides for unique naming
        for i, slide in enumerate(prs.slides):
            slide_num = i + 1
            slide_text_content = ""
            slide_meta = {"source": filename, "slide_number": slide_num, "content_type": "slide_text"}
            slide_urls = set()

            try:
                notes_slide = slide.notes_slide
                if notes_slide and notes_slide.notes_text_frame:
                     notes_text = notes_slide.notes_text_frame.text
                     if notes_text:
                          slide_text_content += f"\n--- Speaker Notes (Slide {slide_num}) ---\n{notes_text}\n"
                          slide_urls.update(find_urls(notes_text))
            except Exception as notes_err:
                 log_message(f"Could not extract speaker notes for slide {slide_num}: {notes_err}", "warning")


            for shape in slide.shapes:
                shape_text = ""
                # Extract text from shapes
                if shape.has_text_frame:
                    try: shape_text = shape.text_frame.text
                    except AttributeError: pass # Ignore shapes without text frames gracefully
                    except Exception as text_err: # Catch other potential errors
                        log_message(f"Error getting text from shape {shape.shape_id} on slide {slide_num}: {text_err}", "debug")
                    if shape_text:
                         slide_text_content += shape_text + "\n"
                         slide_urls.update(find_urls(shape_text))

                # Extract text from tables
                if shape.has_table:
                    table_text = f"\n[Table on Slide {slide_num}:]\n"
                    try:
                        for row in shape.table.rows:
                            row_data = [cell.text_frame.text for cell in row.cells]
                            table_text += "| " + " | ".join(row_data) + " |\n"
                        slide_text_content += table_text
                        slide_urls.update(find_urls(table_text)) # Find URLs in tables too
                    except Exception as table_err:
                         log_message(f"Error reading table on slide {slide_num}: {table_err}", "warning")

                # Handle hyperlinks associated with shapes or text runs
                try:
                    # Shape click action hyperlink
                    if shape.click_action and shape.click_action.hyperlink and shape.click_action.hyperlink.address:
                        slide_urls.add(shape.click_action.hyperlink.address)
                    # Text run hyperlinks
                    if hasattr(shape, 'text_frame'): # Check text runs for hyperlinks
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                if run.hyperlink and run.hyperlink.address:
                                    slide_urls.add(run.hyperlink.address)
                except Exception as link_err:
                    log_message(f"Minor error reading hyperlink on slide {slide_num}: {link_err}", "debug")


                # Handle images (pictures inserted or shape fills)
                if hasattr(shape, 'image'):
                    try:
                        image = shape.image
                        image_bytes = image.blob
                        # Use a more robust unique name
                        img_meta_filename = f"{Path(filename).stem}_slide{slide_num}_img{img_count}_{shape.shape_id}"
                        img_desc_ocr = get_gemini_vision_description_ocr(image_bytes, img_meta_filename)
                        if img_desc_ocr and "No description or text extracted" not in img_desc_ocr:
                             img_meta = {"source": filename, "slide_number": slide_num, "content_type": "embedded_image", "image_ref": f"shape_{shape.shape_id}"}
                             parsed_segments.append((f"[Analysis of Image {img_count} on Slide {slide_num}]\n{img_desc_ocr}", img_meta))
                             img_count += 1
                    except Exception as img_err:
                         log_message(f"Could not process shape as image on slide {slide_num} (ID {shape.shape_id}): {img_err}", "warning")
                # Check for picture shapes explicitly if hasattr('image') fails
                elif shape.shape_type == 13: # MSO_SHAPE_TYPE.PICTURE
                     try:
                          image = shape.image
                          image_bytes = image.blob
                          img_meta_filename = f"{Path(filename).stem}_slide{slide_num}_img{img_count}_{shape.shape_id}"
                          img_desc_ocr = get_gemini_vision_description_ocr(image_bytes, img_meta_filename)
                          if img_desc_ocr and "No description or text extracted" not in img_desc_ocr:
                               img_meta = {"source": filename, "slide_number": slide_num, "content_type": "embedded_image", "image_ref": f"shape_{shape.shape_id}_pic"}
                               parsed_segments.append((f"[Analysis of Picture {img_count} on Slide {slide_num}]\n{img_desc_ocr}", img_meta))
                               img_count += 1
                     except Exception as img_err:
                          log_message(f"Could not process picture shape on slide {slide_num} (ID {shape.shape_id}): {img_err}", "warning")


            # Add combined text for the slide as one segment if not empty
            slide_text_content_stripped = slide_text_content.strip()
            if slide_text_content_stripped:
                 parsed_segments.append((slide_text_content_stripped, slide_meta)) # Add stripped text
            urls.update(slide_urls)

        log_message(f"Finished processing PPTX '{filename}'. Found URLs: {len(urls)}, Images Analyzed: {img_count}")
    except Exception as e:
        log_message(f"Critical error processing PPTX '{filename}': {e}", "error")
        return [], []
    return parsed_segments, list(urls)

# --- Data File Parser (Integrates DataImporter_Gemini) ---
def parse_data_file(file_content, filename):
    """
    Parses CSV, XLSX, or JSON by importing it into a chat-specific SQLite DB.
    Returns a single "marker" text segment indicating success/failure for ChromaDB.
    The actual data interaction happens via VannaDataAnalyzer_Gemini in qa_engine.
    """
    parsed_segments = []
    file_type = filename.split('.')[-1].lower()
    chat_id = st.session_state.get("current_chat_id")
    chat_state = st.session_state.chats.get(chat_id)

    if not chat_id or not chat_state:
        log_message(f"Error parsing data file '{filename}': Could not determine current chat state.", "error")
        base_meta = {"source": filename, "content_type": "data_import_failed", "reason": "Chat state unavailable"}
        marker_text = f"[Data File Import Failed: {filename} - Chat state unavailable]"
        parsed_segments.append((marker_text, base_meta))
        return parsed_segments, []

    chat_db_filename = f"chat_{chat_id}.db"
    chat_db_path = os.path.join(CHAT_DB_DIRECTORY, chat_db_filename)
    api_key = st.session_state.get("google_api_key")
    if not api_key:
         log_message(f"Warning during data file import '{filename}': API Key not configured. Analysis will fail later.", "warning")

    importer = None
    imported_tables = []
    import_success = False
    try:
        log_message(f"Processing Data file '{filename}' for chat '{chat_id}' into DB '{chat_db_path}'.", "info")
        importer = DataImporter_Gemini(db_path=chat_db_path, verbose=True)
        import_kwargs = {'if_exists': 'replace'}

        if file_type == "csv":
            table_name = importer._clean_column_name(Path(filename).stem)
            stream = io.BytesIO(file_content)
            try:
                imported_tables = importer.import_csv_from_stream(
                    stream=stream, table_name=table_name, encoding='utf-8', if_exists=import_kwargs['if_exists']
                 )
            except UnicodeDecodeError:
                log_message(f"UTF-8 decode failed for CSV '{filename}', trying latin1.", "warning")
                stream.seek(0)
                imported_tables = importer.import_csv_from_stream(
                    stream=stream, table_name=table_name, encoding='latin1', if_exists=import_kwargs['if_exists']
                )
            except Exception as csv_err:
                raise RuntimeError(f"CSV stream import failed: {csv_err}")

        elif file_type == "xlsx":
            excel_file_stream = io.BytesIO(file_content)
            try:
                 xls = pd.ExcelFile(excel_file_stream)
                 if not xls.sheet_names: raise RuntimeError("Excel file has no sheets.")
                 sheet_name_to_import = xls.sheet_names[0]
                 table_name = importer._clean_column_name(Path(filename).stem)
                 import_kwargs['sheet_name'] = sheet_name_to_import
                 import_kwargs['table_names'] = {sheet_name_to_import: table_name}
                 imported_tables = importer.import_excel(excel_file_stream, **import_kwargs)
            except Exception as excel_err:
                raise RuntimeError(f"Excel import failed: {excel_err}")

        elif file_type == "json":
             table_name = importer._clean_column_name(Path(filename).stem)
             import_kwargs['table_name'] = table_name
             try:
                 stream = io.BytesIO(file_content)
                 df = pd.read_json(stream)
                 importer._import_dataframe(df, table_name, if_exists=import_kwargs['if_exists'])
                 imported_tables = [table_name]
             except Exception as json_err:
                 raise RuntimeError(f"JSON import/parsing failed: {json_err}")

        if imported_tables:
             import_success = True
             chat_state["chat_db_path"] = chat_db_path
             if "imported_tables" not in chat_state: chat_state["imported_tables"] = []
             chat_state["imported_tables"] = list(set(chat_state["imported_tables"] + imported_tables))

             log_message(f"Successfully imported '{filename}' to table(s) {imported_tables} in DB '{chat_db_path}'.")

             # --- *** FIX THE METADATA HERE *** ---
             # Convert the list of tables to a comma-separated string for ChromaDB
             tables_string = ", ".join(imported_tables)
             base_meta = {
                 "source": filename,
                 "content_type": "data_import_success",
                 "database": chat_db_path,
                 "tables": tables_string # Store as string
             }
             # --- *** END FIX *** ---

             marker_text = f"[Data File Imported: {filename} ({tables_string}) - Ready for querying in chat database.]"
             parsed_segments.append((marker_text, base_meta))
             chat_state["processed_files"][filename] = 'success'

        else:
             log_message(f"Data import process for '{filename}' completed but reported no imported tables (empty file?).", "warning")
             base_meta = {"source": filename, "content_type": "data_import_failed", "reason": "No tables reported by importer"}
             marker_text = f"[Data File Import Failed: {filename} - No tables were imported (file empty or error).]"
             parsed_segments.append((marker_text, base_meta))
             chat_state["processed_files"][filename] = 'failed'

    except Exception as e:
        st.exception(e)
        log_message(f"Error importing data file '{filename}' into database: {e}", "error")
        base_meta = {"source": filename, "content_type": "data_import_failed", "reason": str(e)}
        # Ensure 'tables' key doesn't exist or is a string if adding failure marker metadata
        # base_meta["tables"] = "" # Add empty string if needed
        marker_text = f"[Data File Import Failed: {filename} - Error: {e}]"
        parsed_segments.append((marker_text, base_meta))
        if filename in chat_state.get("processed_files", {}):
             chat_state["processed_files"][filename] = 'failed'
        return parsed_segments, []
    finally:
        if importer:
            importer.close()

    return parsed_segments, []


# --- TXT Parser ---
# ... (keep as is) ...
def parse_txt(file_content, filename):
    text = ""
    urls = set()
    base_meta = {"source": filename, "content_type": "text"}
    parsed_segments = []
    try:
        log_message(f"Processing TXT file '{filename}'.", "info")
        try: text = file_content.decode("utf-8")
        except UnicodeDecodeError: text = file_content.decode("latin1", errors='ignore')

        if text:
             urls.update(find_urls(text))
             parsed_segments.append((text, base_meta))
             log_message(f"Finished processing TXT '{filename}'. Found URLs: {len(urls)}")
        else:
             log_message(f"TXT file '{filename}' is empty.", "warning")

    except Exception as e:
        log_message(f"Error processing TXT file '{filename}': {e}", "error")
        return [], []
    return parsed_segments, list(urls)

# --- Image File Parser ---
# ... (keep as is) ...
def parse_image(file_content, filename):
    text = ""
    base_meta = {"source": filename, "content_type": "image_analysis"}
    parsed_segments = []
    try:
        log_message(f"Processing Image file '{filename}'.", "info")
        img_desc_ocr = get_gemini_vision_description_ocr(file_content, filename)
        if img_desc_ocr and "No description or text extracted" not in img_desc_ocr:
             text = f"[Analysis of Image File: {filename}]\n{img_desc_ocr}\n"
             parsed_segments.append((text, base_meta))
             log_message(f"Finished processing image file: {filename}")
        else:
             log_message(f"Image analysis did not return significant content for {filename}.", "warning")

    except Exception as e:
        log_message(f"Error processing image file '{filename}': {e}", "error")
        return [], []
    return parsed_segments, [] # No URLs expected
# --- END OF FILE file_parsers.py ---